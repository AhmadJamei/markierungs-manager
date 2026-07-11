import os

from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse
from django.contrib.auth.decorators import login_required
from django.db.models import Q
from django.conf import settings
from django.core.cache import cache
from django.http import HttpResponse
from datetime import datetime, timedelta, date
from Vehicle.models import Vehicle
from Contract.models import Contract
from Accounts.models import CustomUser
from openpyxl.styles import Font, PatternFill, Alignment
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from .models import WorkDay, WorkDayWorkerNote, WorkReport, WorkReportImage, WorkReportAudio
from io import BytesIO
import requests as req
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
import requests
import openpyxl


@login_required
def schedule_view(request):
    
    # دریافت هفته از URL یا هفته جاری
    week_str = request.GET.get('week')
    if week_str:
        week_start = datetime.strptime(week_str, '%Y-%m-%d').date()
    else:
        today = datetime.now().date()
        week_start = today - timedelta(days=today.weekday())

    week_end = week_start + timedelta(days=6)

    week_number = week_start.isocalendar().week
    year = week_start.isocalendar().year    
    print("WEEK NUM:", week_number)
    # روزهای هفته
    week_days = [week_start + timedelta(days=i) for i in range(7)]

    is_worker = request.user.groups.filter(name='worker').exists()

    # برنامه هفته
    workdays = WorkDay.objects.filter(
        date__range=[week_start, week_end]
    ).select_related('contract', 'engineer').prefetch_related(
        'workers', 'vehicles',
        'worker_notes', 'worker_notes__worker',
        'reports', 'reports__worker', 'reports__images')


    # قراردادهای این هفته
    contracts_in_week = Contract.objects.filter(
        workdays__date__range=[week_start, week_end]
    ).distinct()    # ساختار جدول
    schedule_data = {}
    for day in week_days:
        schedule_data[day] = workdays.filter(date=day)

    context = {
        'week_days': week_days,
        'schedule_data': schedule_data,
        'week_start': week_start,
        'week_end': week_end,
        'week_number': week_number,
        'prev_week': (week_start - timedelta(days=7)).strftime('%Y-%m-%d'),
        'next_week': (week_start + timedelta(days=7)).strftime('%Y-%m-%d'),
        'contracts': Contract.objects.all(),
        'current_week': week_start.strftime('%Y-%m-%d'),
        'workers': CustomUser.objects.filter(
            groups__name='worker'
        ).exclude(
            employee__status='terminated'
        ),
        'vehicles': Vehicle.objects.all(),  # ← اضافه شد

    }
    return render(request, 'Schedule/schedule.html', context)

@login_required
def add_workday(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        contract_id = data.get('contract_id')
        start_date = data.get('start_date')
        end_date = data.get('end_date', start_date)
        worker_ids = data.get('worker_ids', [])
        vehicle_ids = data.get('vehicle_ids', [])
        note = data.get('note', '')
        shift = data.get('shift', 'day')

        contract = get_object_or_404(Contract, pk=contract_id)
        vehicles = Vehicle.objects.filter(RegisterNumber__in=vehicle_ids)

        start = datetime.strptime(start_date, '%Y-%m-%d').date()
        end = datetime.strptime(end_date, '%Y-%m-%d').date()

        created_count = 0
        current = start
        while current <= end:
            workday, created = WorkDay.objects.get_or_create(
                contract=contract,
                date=current,
                defaults={
                    'engineer': request.user,
                    'note': note,
                    'shift': shift,
                }
            )
            workday.workers.set(worker_ids)
            workday.vehicles.set(vehicles)
            workday.save()
            created_count += 1
            current += timedelta(days=1)

        return JsonResponse({'status': 'ok', 'created': created_count})
    return JsonResponse({'status': 'error'}, status=400)

@login_required
def delete_workday(request, pk):
    workday = get_object_or_404(WorkDay, pk=pk)
    if request.method == 'POST':
        workday.delete()
        return JsonResponse({'status': 'ok'})
    return JsonResponse({'status': 'error'}, status=400)

@login_required
def get_workday_workers(request, pk):
    workday = get_object_or_404(WorkDay, pk=pk)
    workers = list(workday.workers.values('id', 'username', 'first_name', 'last_name'))
    return JsonResponse({'workers': workers})

@login_required
def update_workday(request, pk):
    workday = get_object_or_404(WorkDay, pk=pk)
    if request.method == 'POST':
        data = json.loads(request.body)
        worker_ids = data.get('worker_ids', [])
        vehicle_ids = data.get('vehicle_ids', [])
        note = data.get('note', '')
        shift = data.get('shift', workday.shift)  # ← اضافه شد
        
        workday.workers.set(worker_ids)
        vehicles = Vehicle.objects.filter(RegisterNumber__in=vehicle_ids)
        workday.vehicles.set(vehicles)
        workday.note = note
        workday.shift = shift  # ← اضافه شد
        workday.save()
        return JsonResponse({'status': 'ok'})
    return JsonResponse({'status': 'error'}, status=400)

@login_required
def copy_workday(request, pk):
    workday = get_object_or_404(WorkDay, pk=pk)
    if request.method == 'POST':
        data = json.loads(request.body)
        new_date = data.get('date')

        new_workday, created = WorkDay.objects.get_or_create(
            contract=workday.contract,
            date=new_date,
            defaults={
                'engineer': request.user,
                'note': workday.note
            }
        )
        new_workday.workers.set(workday.workers.all())
        new_workday.vehicles.set(workday.vehicles.all())
        new_workday.save()
        return JsonResponse({'status': 'ok'})
    return JsonResponse({'status': 'error'}, status=400)

@login_required
def get_workday_vehicles(request, pk):
    workday = get_object_or_404(WorkDay, pk=pk)
    vehicles = list(workday.vehicles.values('RegisterNumber', 'Model'))
    return JsonResponse({'vehicles': vehicles})

@login_required
def get_weather(request):
    city = request.GET.get('city')
    date_str = request.GET.get('date')
    shift = request.GET.get('shift', 'day')
    
    if not city:
        return JsonResponse({'status': 'error', 'message': 'No city provided'})
    
    # تعیین ساعت بر اساس shift
    target_hour = 12 if shift == 'day' else 21
    
    # چک cache
    cache_key = f'weather_{city}_{date_str}_{shift}'
    cached = cache.get(cache_key)
    if cached:
        return JsonResponse({'status': 'ok', 'weather': cached})
    
    try:
        # اگه تاریخ امروز یا گذشته بود، current weather بگیر
        from datetime import date as date_type
        today = date_type.today()
        target_date = datetime.strptime(date_str, '%Y-%m-%d').date() if date_str else today
        
        if target_date <= today:
            # Current weather
            url = "https://api.openweathermap.org/data/2.5/weather"
            params = {
                'q': city,
                'appid': settings.OPENWEATHER_API_KEY,
                'units': 'metric',
                'lang': 'de'
            }
            response = requests.get(url, params=params)
            data = response.json()

            if data.get('cod') != 200:
                return JsonResponse({'status': 'error', 'message': 'City not found'})

            weather = {
                'temp': round(data['main']['temp']),
                'feels_like': round(data['main']['feels_like']),
                'description': data['weather'][0]['description'],
                'icon': data['weather'][0]['icon'],
                'wind_speed': data['wind']['speed'],
                'humidity': data['main']['humidity'],
                'rain': 'rain' in data or 'drizzle' in data['weather'][0]['main'].lower(),
            }
        else:
            # Forecast weather
            url = "https://api.openweathermap.org/data/2.5/forecast"
            params = {
                'q': city,
                'appid': settings.OPENWEATHER_API_KEY,
                'units': 'metric',
                'lang': 'de'
            }
            response = requests.get(url, params=params)
            data = response.json()

            if data.get('cod') != '200':
                return JsonResponse({'status': 'error', 'message': 'City not found'})

            # پیدا کردن نزدیک‌ترین پیش‌بینی به تاریخ و ساعت مورد نظر
            target_datetime = f"{date_str} {target_hour:02d}:00:00"
            best_match = None
            for item in data['list']:
                if item['dt_txt'] >= target_datetime:
                    best_match = item
                    break
            
            if not best_match:
                best_match = data['list'][-1]

            weather = {
                'temp': round(best_match['main']['temp']),
                'feels_like': round(best_match['main']['feels_like']),
                'description': best_match['weather'][0]['description'],
                'icon': best_match['weather'][0]['icon'],
                'wind_speed': best_match['wind']['speed'],
                'humidity': best_match['main']['humidity'],
                'rain': 'rain' in best_match or 'rain' in best_match['weather'][0]['main'].lower(),
            }

        # ذخیره در cache برای ۳۰ دقیقه
        cache.set(cache_key, weather, 60 * 30)
        return JsonResponse({'status': 'ok', 'weather': weather})
    
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)})
    
@login_required
def update_status(request, pk):
    workday = get_object_or_404(WorkDay, pk=pk)
    if request.method == 'POST':
        data = json.loads(request.body)
        status = data.get('status')
        cancel_reason = data.get('cancel_reason', '')
        workday.status = status
        workday.cancel_reason = cancel_reason
        workday.save()
        return JsonResponse({'status': 'ok'})
    return JsonResponse({'status': 'error'}, status=400)   

@login_required
def worker_report(request):
    if not request.user.groups.filter(name__in=['manager', 'engineer']).exists():
        return render(request, 'Dashboard/no_access.html')

    # فیلترها
    from_date = request.GET.get('from_date')
    to_date = request.GET.get('to_date')
    worker_id = request.GET.get('worker_id')
    contract_id = request.GET.get('contract_id')

    workdays = WorkDay.objects.all()

    if from_date:
        workdays = workdays.filter(date__gte=from_date)
    if to_date:
        workdays = workdays.filter(date__lte=to_date)
    if contract_id:
        workdays = workdays.filter(contract_id=contract_id)

    # گزارش بر اساس کارمند
    workers = CustomUser.objects.filter(groups__name='worker')
    if worker_id:
        workers = workers.filter(id=worker_id)

    report_data = []
    for worker in workers:
        worker_workdays = workdays.filter(workers=worker)
        report_data.append({
            'worker': worker,
            'total_days': worker_workdays.count(),
            'done_days': worker_workdays.filter(status='done').count(),
            'cancelled_days': worker_workdays.exclude(status__in=['planned', 'done']).count(),
            'planned_days': worker_workdays.filter(status='planned').count(),
            'contracts': worker_workdays.values('contract__IDContract', 'contract__Name').distinct(),
        })

    total_days = sum(item['total_days'] for item in report_data)
    total_done = sum(item['done_days'] for item in report_data)
    total_cancelled = sum(item['cancelled_days'] for item in report_data)
    total_planned = sum(item['planned_days'] for item in report_data)

    # مرتب‌سازی بر اساس تعداد روزهای کاری
    report_data.sort(key=lambda x: x['total_days'], reverse=True)

    context = {
        'report_data': report_data,
        'workers': CustomUser.objects.filter(groups__name='worker'),
        'contracts': Contract.objects.all(),
        'from_date': from_date or '',
        'to_date': to_date or '',
        'selected_worker': worker_id or '',
        'selected_contract': contract_id or '',
        'total_days': total_days,
        'total_done': total_done,
        'total_cancelled': total_cancelled,
        'total_planned': total_planned,
    }
    print(report_data)
    return render(request, 'Schedule/worker_report.html', context)

@login_required
def worker_report_excel(request):
    if not request.user.groups.filter(name__in=['manager', 'engineer']).exists():
        return JsonResponse({'status': 'error'}, status=403)

    from_date = request.GET.get('from_date')
    to_date = request.GET.get('to_date')
    worker_id = request.GET.get('worker_id')
    contract_id = request.GET.get('contract_id')

    workdays = WorkDay.objects.all()
    if from_date:
        workdays = workdays.filter(date__gte=from_date)
    if to_date:
        workdays = workdays.filter(date__lte=to_date)
    if contract_id:
        workdays = workdays.filter(contract_id=contract_id)

    workers = CustomUser.objects.filter(groups__name='worker')
    if worker_id:
        workers = workers.filter(id=worker_id)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Worker Report'

    headers = ['Worker', 'Total Days', 'Done', 'Planned', 'Cancelled', 'Contracts']
    ws.append(headers)

    for cell in ws[1]:
        cell.font = Font(bold=True, color='FFFFFF')
        cell.fill = PatternFill(fill_type='solid', fgColor='0d6efd')
        cell.alignment = Alignment(horizontal='center')

    for worker in workers:
        worker_workdays = workdays.filter(workers=worker)
        contracts = ', '.join(worker_workdays.values_list('contract__IDContract', flat=True).distinct())
        ws.append([
            worker.get_full_name() or worker.username,
            worker_workdays.count(),
            worker_workdays.filter(status='done').count(),
            worker_workdays.filter(status='planned').count(),
            worker_workdays.exclude(status__in=['planned', 'done']).count(),
            contracts,
        ])

    ws.column_dimensions['A'].width = 25
    ws.column_dimensions['B'].width = 15
    ws.column_dimensions['C'].width = 15
    ws.column_dimensions['D'].width = 15
    ws.column_dimensions['E'].width = 15
    ws.column_dimensions['F'].width = 30

    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    response['Content-Disposition'] = 'attachment; filename=worker_report.xlsx'
    wb.save(response)
    return response

@login_required
def worker_report_pdf(request):
    if not request.user.groups.filter(name__in=['manager', 'engineer']).exists():
        return JsonResponse({'status': 'error'}, status=403)

    from_date = request.GET.get('from_date')
    to_date = request.GET.get('to_date')
    worker_id = request.GET.get('worker_id')
    contract_id = request.GET.get('contract_id')

    workdays = WorkDay.objects.all()
    if from_date:
        workdays = workdays.filter(date__gte=from_date)
    if to_date:
        workdays = workdays.filter(date__lte=to_date)
    if contract_id:
        workdays = workdays.filter(contract_id=contract_id)

    workers = CustomUser.objects.filter(groups__name='worker')
    if worker_id:
        workers = workers.filter(id=worker_id)

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename=worker_report.pdf'

    doc = SimpleDocTemplate(
        response,
        pagesize=landscape(A4),
        leftMargin=20, rightMargin=20,
        topMargin=20, bottomMargin=20,
    )
    elements = []

    data = [['Worker', 'Total Days', 'Done', 'Planned', 'Cancelled', 'Contracts']]
    for worker in workers:
        worker_workdays = workdays.filter(workers=worker)
        contracts = ', '.join(worker_workdays.values_list('contract__IDContract', flat=True).distinct())
        data.append([
            worker.get_full_name() or worker.username,
            worker_workdays.count(),
            worker_workdays.filter(status='done').count(),
            worker_workdays.filter(status='planned').count(),
            worker_workdays.exclude(status__in=['planned', 'done']).count(),
            contracts,
        ])

    table = Table(data, repeatRows=1)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0d6efd')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f2f2f2')]),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('PADDING', (0, 0), (-1, -1), 5),
    ]))

    elements.append(table)
    doc.build(elements)
    return response

@login_required
def save_worker_note(request, workday_id, worker_id):
    if request.method == 'POST':
        data = json.loads(request.body)
        note_text = data.get('note', '')
        
        note, created = WorkDayWorkerNote.objects.get_or_create(
            workday_id=workday_id,
            worker_id=worker_id,
        )
        note.note = note_text
        note.save()
        
        return JsonResponse({'status': 'ok'})
    return JsonResponse({'status': 'error'}, status=400)

@login_required
def get_worker_note(request, workday_id, worker_id):
    try:
        note = WorkDayWorkerNote.objects.get(
            workday_id=workday_id,
            worker_id=worker_id
        )
        return JsonResponse({'status': 'ok', 'note': note.note})
    except WorkDayWorkerNote.DoesNotExist:
        return JsonResponse({'status': 'ok', 'note': ''})
    
@login_required
def my_schedule(request):
    week_str = request.GET.get('week')
    if week_str:
        week_start = datetime.strptime(week_str, '%Y-%m-%d').date()
    else:
        today = datetime.now().date()
        week_start = today - timedelta(days=today.weekday())

    week_end = week_start + timedelta(days=6)
    week_number = week_start.isocalendar().week
    week_days = [week_start + timedelta(days=i) for i in range(7)]

    # فقط workday های این کارمند
    workdays = WorkDay.objects.filter(
        date__range=[week_start, week_end],
        workers=request.user
    ).select_related('contract', 'engineer').prefetch_related(
        'workers', 'vehicles', 'worker_notes', 'worker_notes__worker'
    )

    # ساختار جدول
    schedule_data = {}
    for day in week_days:
        schedule_data[day] = workdays.filter(date=day)

    context = {
        'week_days': week_days,
        'schedule_data': schedule_data,
        'week_start': week_start,
        'week_end': week_end,
        'week_number': week_number,
        'prev_week': (week_start - timedelta(days=7)).strftime('%Y-%m-%d'),
        'next_week': (week_start + timedelta(days=7)).strftime('%Y-%m-%d'),
    }
    return render(request, 'Schedule/my_schedule.html', context)

@login_required
def add_report(request, workday_id):
    workday = get_object_or_404(WorkDay, pk=workday_id)
    if request.method == 'POST':
        text = request.POST.get('text', '')
        
        # گرفتن یا ساختن report
        report, created = WorkReport.objects.get_or_create(
            workday=workday,
            worker=request.user,
            defaults={'text': text}
        )
        
        if not created:
            report.text = text
            report.save()
        
        # آپلود تصاویر
        images = request.FILES.getlist('images')
        for image in images:
            WorkReportImage.objects.create(
                report=report,
                image=image,
                caption=request.POST.get('caption', '')
            )
        # آپلود صدا
        audio_file = request.FILES.get('audio')
        if audio_file:
            WorkReportAudio.objects.create(
                report=report,
                audio=audio_file,
            )

        return JsonResponse({'status': 'ok', 'report_id': report.id})
    return JsonResponse({'status': 'error'}, status=400)

@login_required
def get_report(request, workday_id):
    try:
        report = WorkReport.objects.get(
            workday_id=workday_id,
            worker=request.user
        )
        images = [{'id': img.id, 'url': img.image.url, 'caption': img.caption} 
                  for img in report.images.all()]
        audios = [{'id': aud.id, 'url': aud.audio.url} 
                  for aud in report.audios.all()]
        return JsonResponse({
            'status': 'ok',
            'report_id': report.id,
            'text': report.text,
            'status_value': report.status,
            'engineer_note': report.engineer_note,
            'images': images,
            'audios': audios,
        })
    except WorkReport.DoesNotExist:
        return JsonResponse({'status': 'ok', 'text': '', 'images': [], 'audios': []})

@login_required
def delete_report_image(request, image_id):
    image = get_object_or_404(WorkReportImage, pk=image_id)
    if image.report.worker == request.user:
        image.delete()
        return JsonResponse({'status': 'ok'})
    return JsonResponse({'status': 'error'}, status=403)

@login_required
def review_report(request, report_id):
    if not request.user.groups.filter(name__in=['manager', 'engineer']).exists():
        return JsonResponse({'status': 'error'}, status=403)
    
    report = get_object_or_404(WorkReport, pk=report_id)
    if request.method == 'POST':
        data = json.loads(request.body)
        report.status = data.get('status')
        report.engineer_note = data.get('engineer_note', '')
        report.save()
        return JsonResponse({'status': 'ok'})
    return JsonResponse({'status': 'error'}, status=400)

@login_required
def contract_gallery(request):
    from Contract.models import Contract
    from django.db.models import Count
    
    contracts = Contract.objects.filter(
        workdays__reports__images__isnull=False
    ).annotate(
        image_count=Count('workdays__reports__images', distinct=True)
    ).prefetch_related(
        'workdays__reports__images',
        'workdays__reports__worker',
    ).distinct()
    selected_contract = request.GET.get('contract')  # ← اضافه شد

    context = {
        'contracts': contracts,
        'selected_contract': selected_contract,
    }
    return render(request, 'Schedule/contract_gallery.html', context)

@login_required
def workday_detail(request, pk):
    workday = get_object_or_404(WorkDay, pk=pk)
    workday = WorkDay.objects.filter(pk=pk).select_related(
        'contract', 'engineer'
    ).prefetch_related(
        'workers',
        'vehicles',
        'worker_notes__worker',
        'reports__worker',
        'reports__images',
        'contract__materials__material',
        'reports__audios',  #*  اضافه شد

    ).first()

    context = {
        'workday': workday,
    }
    return render(request, 'Schedule/workday_detail.html', context)

def duplicate_slide(prs, source_slide):
    import copy
    from pptx.oxml.ns import qn
    
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    # کپی spTree
    source_spTree = source_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    new_spTree = new_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    
    for child in list(new_spTree):
        new_spTree.remove(child)
    
    for child in source_spTree:
        new_spTree.append(copy.deepcopy(child))
    
    # کپی background
    source_bg = source_slide._element.find(qn('p:cSld')).find(qn('p:bg'))
    if source_bg is not None:
        new_cSld = new_slide._element.find(qn('p:cSld'))
        existing_bg = new_cSld.find(qn('p:bg'))
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        new_cSld.insert(0, copy.deepcopy(source_bg))
    
    # کپی transition و timing
    for tag in [qn('p:transition'), qn('p:timing')]:
        el = source_slide._element.find(tag)
        if el is not None:
            existing = new_slide._element.find(tag)
            if existing is not None:
                new_slide._element.remove(existing)
            new_slide._element.append(copy.deepcopy(el))
    
    return new_slide

@login_required
def generate_weekly_pptx(request):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from io import BytesIO
    import copy
    import os

    # گرفتن هفته
    week_str = request.GET.get('week')
    if week_str:
        week_start = datetime.strptime(week_str, '%Y-%m-%d').date()
    else:
        today = datetime.now().date()
        week_start = today - timedelta(days=today.weekday())

    week_end = week_start + timedelta(days=6)
    week_number = week_start.isocalendar()[1]

    # گرفتن workday ها
    workdays = WorkDay.objects.filter(
        date__range=[week_start, week_end]
    ).select_related('contract').prefetch_related(
        'workers', 'workers__employee'
    ).order_by('date')

    from Settings.models import CompanySettings
    company = CompanySettings.objects.first()

    DAY_NAMES = {0:'Monday', 1:'Tuesday', 2:'Wednesday', 3:'Thursday', 4:'Friday', 5:'Saturday', 6:'Sunday'}

    template_path = os.path.join(settings.BASE_DIR, 'Schedule', 'static', 'Schedule', 'templates', 'weekly_template.pptx')
    default_photo = os.path.join(settings.BASE_DIR, 'static', 'assets', 'img', 'Default Picture.jpg')

    def replace_text_in_slide(slide, replacements):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = ''.join([run.text for run in para.runs])
                    needs_replace = any(key in full_text for key in replacements)
                    if needs_replace:
                        new_text = full_text
                        for key, val in replacements.items():
                            new_text = new_text.replace(key, str(val))
                        if para.runs:
                            para.runs[0].text = new_text
                            for run in para.runs[1:]:
                                run.text = ''



    def replace_picture_in_slide(slide, placeholder, image_path):
        from pptx.oxml.ns import qn
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import copy
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt = shape._element.nvPicPr.cNvPr.get('descr', '')
                if alt == placeholder:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    # ذخیره spPr (شامل ellipse و crop)
                    spPr = copy.deepcopy(shape._element.find(qn('p:spPr')))
                    
                    # حذف عکس قدیمی
                    sp = shape._element
                    sp.getparent().remove(sp)
                    
                    try:
                        # اضافه کردن عکس جدید
                        new_pic = slide.shapes.add_picture(image_path, left, top, width, height)
                        
                        # جایگزینی spPr با نسخه اصلی (که ellipse داره)
                        if spPr is not None:
                            old_spPr = new_pic._element.find(qn('p:spPr'))
                            if old_spPr is not None:
                                new_pic._element.remove(old_spPr)
                            new_pic._element.append(spPr)
                    except:
                        pass
                    return


    def get_weather_text(workday):
        try:
            city = workday.contract.City
            if not city:
                return ''
            from datetime import date as date_type
            today = date_type.today()
            if workday.date <= today:
                url = "https://api.openweathermap.org/data/2.5/weather"
                params = {'q': city, 'appid': settings.OPENWEATHER_API_KEY, 'units': 'metric', 'lang': 'de'}
            else:
                url = "https://api.openweathermap.org/data/2.5/forecast"
                params = {'q': city, 'appid': settings.OPENWEATHER_API_KEY, 'units': 'metric', 'lang': 'de'}
            r = requests.get(url, params=params, timeout=3)
            data = r.json()
            if workday.date <= today:
                temp = round(data['main']['temp'])
                wind = round(data['wind']['speed'])
                desc = data['weather'][0]['description']
            else:
                date_str = workday.date.strftime('%Y-%m-%d')
                target_hour = 12 if workday.shift == 'day' else 21
                target_dt = f"{date_str} {target_hour:02d}:00:00"
                best = data['list'][-1]
                for item in data['list']:
                    if item['dt_txt'] >= target_dt:
                        best = item
                        break
                temp = round(best['main']['temp'])
                wind = round(best['wind']['speed'])
                desc = best['weather'][0]['description']
            return f'{temp}°C  |  💨 {wind}m/s  |  {desc}'
        except:
            return ''

    # لود template
    prs = Presentation(template_path)

    # اسلاید ۱ - عنوان
    title_slide = prs.slides[0]
    replace_text_in_slide(title_slide, {
        '{{COMPANY_NAME}}': company.name if company else '',
        '{{WEEK_NUMBER}}': f'KW {week_number}',
    })

    # نگه داشتن مرجع template قبل از حذف
    template_slide_ref = prs.slides[1]

    # حذف اسلاید template
    slide_id_list = prs.slides._sldIdLst
    template_rId = slide_id_list[1].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(template_rId)
    del slide_id_list[1]

    # اضافه کردن اسلاید برای هر workday
    for workday in workdays:
        workers = list(workday.workers.all())
        weather = get_weather_text(workday)

        replacements = {
            '{{WEEK_NUMBER}}': f'KW {week_number}',
            '{{DAY}}': DAY_NAMES[workday.date.weekday()],
            '{{DATE}}': workday.date.strftime('%d.%m.%Y'),
            '{{CONTRACT_ID}}': workday.contract.IDContract or '',
            '{{CONTRACT_NAME}}': workday.contract.Name or '',
            '{{ADDRESS}}': workday.contract.Address or '',
            '{{WEATHER}}': weather,
        }
        for i in range(1, 7):
            if i <= len(workers):
                replacements[f'{{{{WORKER_{i}_NAME}}}}'] = workers[i-1].get_full_name() or workers[i-1].username
            else:
                replacements[f'{{{{WORKER_{i}_NAME}}}}'] = ''

        # اضافه کردن اسلاید جدید
        new_slide = duplicate_slide(prs, template_slide_ref)

        # جایگزینی متن
        replace_text_in_slide(new_slide, replacements)

        # جایگزینی عکس‌ها
        for i, worker in enumerate(workers[:6], 1):
            photo_path = default_photo
            try:
                if hasattr(worker, 'employee') and worker.employee.photo:
                    photo_path = worker.employee.photo.path
            except:
                pass
            replace_picture_in_slide(new_slide, f'{{{{WORKER_{i}_PHOTO}}}}', photo_path)

        for i in range(len(workers)+1, 7):
            remove_picture_from_slide(new_slide, f'{{{{WORKER_{i}_PHOTO}}}}')

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    response = HttpResponse(
        output.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = f'attachment; filename=weekly_plan_kw{week_number}.pptx'
    return response
def remove_picture_from_slide(slide, placeholder):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            alt = shape._element.nvPicPr.cNvPr.get('descr', '')
            if alt == placeholder:
                shape._element.getparent().remove(shape._element)
                return